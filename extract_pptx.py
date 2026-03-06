import pptx
import sys

def extract_text(file_path):
    try:
        prs = pptx.Presentation(file_path)
        text_content = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text)
        with open("pptx_content.txt", "w", encoding="utf-8") as f:
            f.write("\n\n".join(text_content))
        print("Success")
    except Exception as e:
        print(f"Error: {e}")

if __name__ == "__main__":
    if len(sys.argv) > 1:
        extract_text(sys.argv[1])
    else:
        print("Provide file path.")
